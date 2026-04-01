"""
data_parser.py
--------------
Data Parsing tool for the LBO multi-agent system.

Extracts financial data from uploaded documents (.xlsx, .pptx, .pdf, .docx,
images) and returns a strictly standardized JSON schema that downstream agents
can write directly into the LBO Excel model.

Usage (Python API):
    from data_parser import parse_document

    result = parse_document("CIM_AcmeCorp.pdf", agent_id="extraction_agent")
    print(result["metadata"]["company_name"])
    print(result["historical"][0]["ebitda"]["value"])
    print(result["warnings"])

Environment variables:
    ANTHROPIC_API_KEY   Required for PDF/PPTX/image/DOCX parsing via Claude vision.
"""

import base64
import json
import os
import re
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional, Union

# ---------------------------------------------------------------------------
# Lazy imports — only loaded when the relevant file type is encountered
# ---------------------------------------------------------------------------
def _import_openpyxl():
    import openpyxl
    return openpyxl

def _import_pdfplumber():
    import pdfplumber
    return pdfplumber

def _import_pptx():
    from pptx import Presentation
    return Presentation

def _import_docx():
    import docx
    return docx

def _import_anthropic():
    import anthropic
    return anthropic


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS = {".xlsx", ".pptx", ".pdf", ".docx", ".png", ".jpg", ".jpeg", ".webp"}

_CLAUDE_MODEL = "claude-opus-4-6"

_VISION_SYSTEM_PROMPT = """You are a financial data extraction specialist for LBO models.
Your job is to extract financial data from documents and return it in a strict JSON schema.
Be precise with numbers. Preserve the exact values shown — do not round or estimate.
If a value is not present, use null. Never fabricate numbers.
Always identify the currency and units (e.g. USD millions, EUR thousands)."""

_EXTRACTION_PROMPT = """Extract ALL financial data from this document and return ONLY a JSON object
matching this exact schema. Do not include any explanation or markdown — just the raw JSON.

Schema:
{
  "company_name": "string or null",
  "fiscal_year_end": "MM-DD or null (e.g. '12-31')",
  "currency": "USD/EUR/GBP/etc or null",
  "units": "millions/thousands/billions or null",
  "historical": [
    {
      "year": integer,
      "revenue": number or null,
      "gross_profit": number or null,
      "gross_margin_pct": number or null,
      "ebitda": number or null,
      "ebitda_margin_pct": number or null,
      "da": number or null,
      "ebit": number or null,
      "net_income": number or null,
      "eps": number or null,
      "total_assets": number or null,
      "total_equity": number or null,
      "total_debt": number or null,
      "cash": number or null,
      "net_debt": number or null,
      "capex": number or null,
      "free_cash_flow": number or null
    }
  ],
  "projected": [
    { same structure as historical }
  ]
}

Rules:
- historical: years with label "actual", "A", or no label (past years)
- projected: years with label "estimate", "E", "forecast", "F", "budget", "LTM", or future years
- All monetary values in the same units as stated in the document
- Percentages as decimals (e.g. 25.3% → 25.3, NOT 0.253)
- capex should be positive (absolute value)
- net_debt = total_debt - cash
- If gross_margin_pct is not shown but gross_profit and revenue are, compute it
- If ebitda_margin_pct is not shown but ebitda and revenue are, compute it
- Return null for any field not found — never guess
"""


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def parse_document(
    file_path: Union[str, Path],
    agent_id: str = "data_parser",
    api_key: Optional[str] = None,
) -> dict:
    """
    Parse a financial document and return standardized LBO JSON schema.

    Parameters
    ----------
    file_path : str | Path
        Path to the document (.xlsx, .pptx, .pdf, .docx, .png, .jpg, .jpeg, .webp).
    agent_id : str
        Identifier for the calling agent (recorded in metadata).
    api_key : str, optional
        Anthropic API key. Falls back to ANTHROPIC_API_KEY env var.

    Returns
    -------
    dict with keys:
        metadata    — company info, source file, extraction timestamp
        historical  — list of annual historical data dicts (each field wrapped in {value, confidence})
        projected   — list of annual projected data dicts (same structure)
        warnings    — list of data integrity warning strings
        raw_text    — raw extracted text/tables (for debugging)
    """
    file_path = Path(file_path)
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: {ext}. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    # Route to extractor
    if ext == ".xlsx":
        raw = _extract_xlsx(file_path)
    elif ext == ".pptx":
        raw = _extract_pptx(file_path, api_key)
    elif ext == ".pdf":
        raw = _extract_pdf(file_path, api_key)
    elif ext == ".docx":
        raw = _extract_docx(file_path, api_key)
    elif ext in {".png", ".jpg", ".jpeg", ".webp"}:
        raw = _extract_image(file_path, api_key)
    else:
        raise ValueError(f"Unhandled extension: {ext}")

    # Normalize raw AI output into final schema
    result = _build_output(raw, file_path, agent_id)
    return result


# ---------------------------------------------------------------------------
# Extractors
# ---------------------------------------------------------------------------

def _extract_xlsx(file_path: Path) -> dict:
    """Deterministic extraction from Excel using openpyxl."""
    openpyxl = _import_openpyxl()
    wb = openpyxl.load_workbook(file_path, data_only=True)

    all_text_blocks = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
        if rows:
            all_text_blocks.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))

    raw_text = "\n\n".join(all_text_blocks)

    # Use AI to map the tabular text to the schema
    return _call_claude_text(raw_text)


def _extract_pdf(file_path: Path, api_key: Optional[str]) -> dict:
    """Extract from PDF: try pdfplumber first, fall back to vision for image pages."""
    pdfplumber = _import_pdfplumber()

    text_pages = []
    image_pages = []  # page indices that had no text

    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            tables = page.extract_tables() or []
            page_content = text.strip()
            for table in tables:
                for row in table:
                    if row:
                        page_content += "\n" + "\t".join(str(c or "") for c in row)
            if page_content.strip():
                text_pages.append(f"--- Page {i+1} ---\n{page_content}")
            else:
                image_pages.append(i)

    raw_text = "\n\n".join(text_pages)

    # If there are image-only pages, send the whole PDF as vision
    if image_pages or not raw_text.strip():
        return _call_claude_vision_file(file_path, api_key, media_type="application/pdf")

    return _call_claude_text(raw_text)


def _extract_pptx(file_path: Path, api_key: Optional[str]) -> dict:
    """Extract text from PPTX slides, then AI-normalize."""
    Presentation = _import_pptx()
    prs = Presentation(str(file_path))

    slide_texts = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = "\t".join(
                        cell.text.strip() for cell in row.cells
                    )
                    if row_text.strip():
                        parts.append(row_text)
        if parts:
            slide_texts.append(f"--- Slide {i+1} ---\n" + "\n".join(parts))

    raw_text = "\n\n".join(slide_texts)

    # If text extraction yielded very little, fall back to vision
    if len(raw_text.strip()) < 200:
        return _call_claude_vision_file(file_path, api_key,
                                        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    return _call_claude_text(raw_text)


def _extract_docx(file_path: Path, api_key: Optional[str]) -> dict:
    """Extract text and tables from a Word document."""
    docx = _import_docx()
    doc = docx.Document(str(file_path))

    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    raw_text = "\n".join(parts)

    if len(raw_text.strip()) < 100:
        return _call_claude_vision_file(file_path, api_key,
                                        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    return _call_claude_text(raw_text)


def _extract_image(file_path: Path, api_key: Optional[str]) -> dict:
    """Send image directly to Claude vision."""
    ext_to_media = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
    }
    media_type = ext_to_media[file_path.suffix.lower()]
    return _call_claude_vision_file(file_path, api_key, media_type=media_type)


# ---------------------------------------------------------------------------
# Claude API calls
# ---------------------------------------------------------------------------

def _get_client(api_key: Optional[str]):
    anthropic = _import_anthropic()
    key = api_key or os.environ.get("ANTHROPIC_API_KEY")
    if not key:
        raise EnvironmentError(
            "ANTHROPIC_API_KEY environment variable is not set. "
            "Set it or pass api_key= to parse_document()."
        )
    return anthropic.Anthropic(api_key=key)


def _call_claude_text(raw_text: str) -> dict:
    """Send extracted text to Claude and ask it to normalize to schema."""
    client = _get_client(None)
    prompt = _EXTRACTION_PROMPT + f"\n\nDocument content:\n\n{raw_text}"

    message = client.messages.create(
        model=_CLAUDE_MODEL,
        max_tokens=4096,
        system=_VISION_SYSTEM_PROMPT,
        messages=[{"role": "user", "content": prompt}],
    )
    response_text = message.content[0].text.strip()
    return {"raw_text": raw_text, "ai_response": response_text}


def _call_claude_vision_file(
    file_path: Path,
    api_key: Optional[str],
    media_type: str,
) -> dict:
    """Encode file as base64 and send to Claude vision."""
    client = _get_client(api_key)

    file_bytes = file_path.read_bytes()
    b64 = base64.standard_b64encode(file_bytes).decode("utf-8")

    # PDFs and Office docs use document type; images use image type
    if media_type.startswith("image/"):
        content_block = {
            "type": "image",
            "source": {"type": "base64", "media_type": media_type, "data": b64},
        }
    else:
        content_block = {
            "type": "document",
            "source": {"type": "base64", "media_type": media_type, "data": b64},
        }

    message = client.messages.create(
        model=_CLAUDE_MODEL,
        max_tokens=4096,
        system=_VISION_SYSTEM_PROMPT,
        messages=[{
            "role": "user",
            "content": [
                content_block,
                {"type": "text", "text": _EXTRACTION_PROMPT},
            ],
        }],
    )
    response_text = message.content[0].text.strip()
    return {"raw_text": None, "ai_response": response_text}


# ---------------------------------------------------------------------------
# Output builder
# ---------------------------------------------------------------------------

_NUMERIC_FIELDS = [
    "revenue", "gross_profit", "gross_margin_pct", "ebitda", "ebitda_margin_pct",
    "da", "ebit", "net_income", "eps",
    "total_assets", "total_equity", "total_debt", "cash", "net_debt",
    "capex", "free_cash_flow",
]

_STRING_FIELDS = ["company_name", "fiscal_year_end", "currency", "units"]


def _parse_ai_json(ai_response: str) -> dict:
    """Parse the JSON block from Claude's response."""
    # Strip markdown code fences if present
    text = re.sub(r"^```(?:json)?\s*", "", ai_response.strip(), flags=re.MULTILINE)
    text = re.sub(r"```\s*$", "", text.strip(), flags=re.MULTILINE)
    try:
        return json.loads(text.strip())
    except json.JSONDecodeError as e:
        raise ValueError(f"Claude returned invalid JSON: {e}\n\nResponse:\n{ai_response[:500]}")


def _wrap_field(value: Any) -> dict:
    """Wrap a raw value with a confidence score."""
    if value is None:
        return {"value": None, "confidence": "not_found"}
    if isinstance(value, (int, float)):
        return {"value": value, "confidence": "high"}
    if isinstance(value, str) and value.strip():
        return {"value": value, "confidence": "high"}
    return {"value": None, "confidence": "not_found"}


def _build_period(raw_period: dict) -> dict:
    """Convert a raw extracted period dict into the wrapped schema."""
    out = {"year": raw_period.get("year")}
    for field in _NUMERIC_FIELDS:
        out[field] = _wrap_field(raw_period.get(field))
    return out


def _validate(historical: list, projected: list) -> list[str]:
    """Run data integrity checks and return a list of warning strings."""
    warnings = []
    all_periods = [("historical", p) for p in historical] + [("projected", p) for p in projected]

    for period_type, period in all_periods:
        year = period.get("year", "?")
        label = f"{period_type} {year}"

        rev = period.get("revenue", {}).get("value")
        ebitda = period.get("ebitda", {}).get("value")
        gross_profit = period.get("gross_profit", {}).get("value")
        ebit = period.get("ebit", {}).get("value")
        net_income = period.get("net_income", {}).get("value")
        total_debt = period.get("total_debt", {}).get("value")
        cash = period.get("cash", {}).get("value")
        net_debt = period.get("net_debt", {}).get("value")
        capex = period.get("capex", {}).get("value")
        ebitda_margin = period.get("ebitda_margin_pct", {}).get("value")
        gross_margin = period.get("gross_margin_pct", {}).get("value")

        if rev is not None and ebitda is not None and ebitda > rev:
            warnings.append(f"[{label}] EBITDA ({ebitda}) exceeds Revenue ({rev}) — likely a units mismatch.")

        if rev is not None and gross_profit is not None and gross_profit > rev:
            warnings.append(f"[{label}] Gross Profit ({gross_profit}) exceeds Revenue ({rev}).")

        if ebitda is not None and ebit is not None and ebit > ebitda:
            warnings.append(f"[{label}] EBIT ({ebit}) exceeds EBITDA ({ebitda}) — D&A would be negative.")

        if ebitda is not None and net_income is not None and net_income > ebitda:
            warnings.append(f"[{label}] Net Income ({net_income}) exceeds EBITDA ({ebitda}) — unexpected.")

        if total_debt is not None and cash is not None and net_debt is not None:
            computed_net_debt = round(total_debt - cash, 4)
            if abs(computed_net_debt - net_debt) > 1.0:
                warnings.append(
                    f"[{label}] Net Debt mismatch: total_debt ({total_debt}) - cash ({cash}) = "
                    f"{computed_net_debt}, but net_debt is {net_debt}."
                )

        if ebitda_margin is not None and rev is not None and ebitda is not None:
            computed_margin = round((ebitda / rev) * 100, 2)
            if abs(computed_margin - ebitda_margin) > 2.0:
                warnings.append(
                    f"[{label}] EBITDA margin mismatch: computed {computed_margin:.1f}% "
                    f"vs stated {ebitda_margin:.1f}%."
                )

        if gross_margin is not None and rev is not None and gross_profit is not None:
            computed_gm = round((gross_profit / rev) * 100, 2)
            if abs(computed_gm - gross_margin) > 2.0:
                warnings.append(
                    f"[{label}] Gross margin mismatch: computed {computed_gm:.1f}% "
                    f"vs stated {gross_margin:.1f}%."
                )

        if capex is not None and capex < 0:
            warnings.append(f"[{label}] Capex is negative ({capex}). Expected a positive absolute value.")

    return warnings


def _build_output(raw: dict, file_path: Path, agent_id: str) -> dict:
    """Parse Claude's JSON response and assemble the final standardized output."""
    ai_response = raw.get("ai_response", "")
    raw_text = raw.get("raw_text")

    extracted = _parse_ai_json(ai_response)

    # Build metadata
    metadata = {
        "company_name": extracted.get("company_name"),
        "fiscal_year_end": extracted.get("fiscal_year_end"),
        "currency": extracted.get("currency"),
        "units": extracted.get("units"),
        "source_file": str(file_path),
        "extracted_at": datetime.now(timezone.utc).isoformat(),
        "agent_id": agent_id,
        "model": _CLAUDE_MODEL,
    }

    # Build historical and projected periods
    historical = [_build_period(p) for p in extracted.get("historical", [])]
    projected = [_build_period(p) for p in extracted.get("projected", [])]

    # Validate
    warnings = _validate(historical, projected)

    return {
        "metadata": metadata,
        "historical": historical,
        "projected": projected,
        "warnings": warnings,
        "raw_text": raw_text,
    }


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python data_parser.py <file_path> [agent_id]")
        sys.exit(1)

    file_arg = sys.argv[1]
    agent_arg = sys.argv[2] if len(sys.argv) > 2 else "cli"

    output = parse_document(file_arg, agent_id=agent_arg)
    print(json.dumps(output, indent=2, default=str))
